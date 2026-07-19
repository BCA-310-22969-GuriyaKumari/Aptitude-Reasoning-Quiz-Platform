#!/usr/bin/env python3
"""
Create a comprehensive, professional PowerPoint presentation for the Quiz Platform project
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)  # Professional blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_shape.line.color.rgb = RGBColor(41, 128, 185)
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.margin_left = Inches(0.3)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(6.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    add_title_slide(prs, 
        "Aptitude & Reasoning Quiz Platform",
        "A Web-Based Interactive Quiz Application\nDesigned by: Guriya Kumari | BCA 6th Semester")
    
    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "🎯 Objective: Create a user-friendly, web-based quiz platform for aptitude and reasoning assessments",
        "✨ Key Features:",
        "   • No login required - simple and quick access",
        "   • Automatic evaluation system",
        "   • Instant result generation",
        "   • Centralized admin management panel",
        "   • Real-time performance tracking"
    ])
    
    # Slide 3: Why This Platform?
    add_content_slide(prs, "Why This Platform?", [
        "❌ Problems with Traditional Tests:",
        "   • Time-consuming and paper-based",
        "   • Manual evaluation leading to errors",
        "   • No instant feedback to students",
        "   • Limited accessibility",
        "",
        "✅ Our Solution:",
        "   • Digital platform eliminates paper",
        "   • Automatic instant evaluation",
        "   • Immediate scorecard generation",
        "   • Accessible from anywhere anytime"
    ])
    
    # Slide 4: Technology Stack
    add_content_slide(prs, "Technology Stack", [
        "Frontend: HTML5, CSS3, Bootstrap (Responsive Design)",
        "Backend: Spring Boot Framework",
        "Database: MySQL (Relational Database)",
        "ORM: Spring Data JPA (Seamless Database Interaction)",
        "IDE: Apache NetBeans 22",
        "Tools: MySQL Workbench (Database Management)",
        "",
        "Hardware: Intel i5 Processor, 4GB RAM, 2GB Storage, Windows OS"
    ])
    
    # Slide 5: System Architecture
    add_content_slide(prs, "System Architecture", [
        "User Flow:",
        "   Student → Home Page → Subject Selection → Quiz Start",
        "   → Question Presentation → Answer Submission → Auto Evaluation",
        "   → Scorecard Display & Results Storage",
        "",
        "Admin Flow:",
        "   Admin Login → Dashboard → Question Management",
        "   → Add/Edit/Delete Questions → View Results & Analytics"
    ])
    
    # Slide 6: Database Design - Questions Table
    add_content_slide(prs, "Database Design: Questions Table", [
        "Stores all quiz questions and their metadata:",
        "",
        "Columns:",
        "   • Q_id: Unique question identifier (Primary Key)",
        "   • Title: Question text",
        "   • optionA, optionB, optionC, optionD: Multiple choice options",
        "   • ans: Correct answer (1-4 indicating option)",
        "   • category: Subject/category of the question",
        "",
        "Purpose: Central repository for all quiz content managed by admin"
    ])
    
    # Slide 7: Database Design - Results Table
    add_content_slide(prs, "Database Design: Results Table", [
        "Stores student quiz results and performance metrics:",
        "",
        "Columns:",
        "   • Id: Unique result record identifier (Primary Key)",
        "   • username: Student name (entered without login)",
        "   • score: Number of correct answers",
        "   • category: Quiz subject attempted",
        "   • timestamp: When quiz was taken",
        "   • total_questions: Questions in the quiz",
        "",
        "Purpose: Permanent record of all quiz attempts for tracking and analytics"
    ])
    
    # Slide 8: User Interface - Student View
    add_content_slide(prs, "User Interface: Student Experience", [
        "Home Page: Welcome with platform introduction",
        "",
        "Student Entry: Simple name entry (no login required)",
        "",
        "Subject Selection: Choose aptitude category",
        "",
        "Quiz Page: Clean interface with question, options, timer",
        "",
        "Result Page: Instant scorecard with:",
        "   • Questions attempted",
        "   • Correct answers",
        "   • Percentage score",
        "   • Category performance"
    ])
    
    # Slide 9: Admin Panel - Overview
    add_content_slide(prs, "Admin Panel: Dashboard Overview", [
        "Secure Authentication: Admin login with credentials",
        "",
        "Dashboard Metrics:",
        "   • Total questions available",
        "   • Total students attempted",
        "   • Total quiz results recorded",
        "   • Recent student performance",
        "",
        "Quick Access to:",
        "   • Question Management",
        "   • Results & Analytics",
        "   • Student Performance Tracking"
    ])
    
    # Slide 10: Admin Panel - Question Management
    add_content_slide(prs, "Admin Panel: Question Management", [
        "Add Questions: Create new questions with 4 options and answer",
        "",
        "Edit Questions: Modify existing questions and answers",
        "",
        "Delete Questions: Remove outdated or incorrect questions",
        "",
        "Category Management: Organize questions by subject",
        "",
        "Bulk Operations: Manage multiple questions efficiently",
        "",
        "Database Sync: All changes immediately reflect in student quiz"
    ])
    
    # Slide 11: Admin Panel - Results & Analytics
    add_content_slide(prs, "Admin Panel: Results & Analytics", [
        "View All Results: Complete database of all quiz attempts",
        "",
        "Student Performance: Individual and aggregate statistics",
        "",
        "Category Analysis: Performance breakdown by subject",
        "",
        "Time Tracking: When students took quizzes",
        "",
        "Success Rates: Percentage of correct answers per category",
        "",
        "Reporting: Export and analyze student performance data"
    ])
    
    # Slide 12: Admin Panel - Security & Access Control
    add_content_slide(prs, "Admin Panel: Security Features", [
        "Authentication: Secure login with admin username and password",
        "",
        "Session Management: Session validation and timeout",
        "",
        "Access Control: Only authenticated admins can modify questions",
        "",
        "Operation Logging: All admin actions tracked",
        "",
        "Password Protection: Hardcoded secure credentials",
        "",
        "Cache Prevention: Headers prevent unauthorized page caching"
    ])
    
    # Slide 13: Data Flow Diagram
    add_content_slide(prs, "Data Flow Diagram", [
        "Student Path:",
        "   Student Input → Quiz Attempt → Answer Storage",
        "   → Auto Evaluation → Result Calculation → Display & Save",
        "",
        "Admin Path:",
        "   Admin Login → Question Management",
        "   → Database Update → Student Quiz Access",
        "",
        "Bidirectional:",
        "   Results Table ← Admin Review & Analytics",
        "   Questions Table ← Admin CRUD Operations"
    ])
    
    # Slide 14: Entity Relationship Diagram
    add_content_slide(prs, "Entity Relationship Model", [
        "Question Entity:",
        "   Q_id (PK) → Title, Options A-D, Answer, Category",
        "",
        "Result Entity:",
        "   Id (PK) → Username, Score, Category, Timestamp",
        "",
        "Relationships:",
        "   • Questions → Used in Quizzes",
        "   • Results → Generated from Quizzes",
        "   • Admin → Manages Questions & Reviews Results",
        "",
        "Admin Role: Acts as management layer for both entities"
    ])
    
    # Slide 15: Key Features Summary
    add_content_slide(prs, "Key Features", [
        "✅ No Login Barrier: Students enter name and start instantly",
        "",
        "✅ Automatic Evaluation: Real-time scoring and result generation",
        "",
        "✅ Category-Based: Organize quizzes by subject",
        "",
        "✅ Admin Control: Manage content without coding",
        "",
        "✅ Secure Admin Panel: Protected access to sensitive operations",
        "",
        "✅ Instant Feedback: Students get scorecard immediately",
        "",
        "✅ Permanent Records: All results stored in database"
    ])
    
    # Slide 16: Future Enhancements
    add_content_slide(prs, "Future Scope & Enhancements", [
        "📱 User Registration: Email-based profiles with history",
        "",
        "📊 Negative Marking: Implement competitive exam rules",
        "",
        "🎓 Difficulty Levels: Easy, Medium, Hard question sets",
        "",
        "📈 Advanced Analytics: Detailed performance reports",
        "",
        "🔐 Role-Based Access: Multiple admin roles",
        "",
        "📱 Mobile App: Native mobile quiz application",
        "",
        "🌐 Multi-language Support: International accessibility"
    ])
    
    # Slide 17: Benefits & Impact
    add_content_slide(prs, "Benefits & Impact", [
        "For Students:",
        "   • Quick, hassle-free quiz participation",
        "   • Instant performance feedback",
        "   • Self-assessment and skill evaluation",
        "",
        "For Educators:",
        "   • Centralized question management",
        "   • Comprehensive performance analytics",
        "   • Easy content updates without coding",
        "",
        "For Institution:",
        "   • Digitalized assessment process",
        "   • Reduced paper usage and manual work",
        "   • Scalable platform for large groups"
    ])
    
    # Slide 18: Conclusion
    add_content_slide(prs, "Conclusion", [
        "The Aptitude & Reasoning Quiz Platform successfully combines:",
        "",
        "✨ User-Friendly Interface for students",
        "🛡️ Secure Admin Management for educators",
        "📊 Reliable Database for data persistence",
        "⚡ Real-Time Evaluation for instant feedback",
        "",
        "This platform demonstrates modern web development practices",
        "and provides a practical solution for educational institutions",
        "to conduct assessments efficiently and effectively."
    ])
    
    # Slide 19: Thank You
    add_title_slide(prs,
        "Thank You!",
        "Aptitude & Reasoning Quiz Platform\n\nFor questions or demonstrations, please contact.\n\nDesigned & Developed: Guriya Kumari | BCA 6th Semester")
    
    return prs

if __name__ == "__main__":
    import os
    os.chdir(r'e:\BCA Final Year Project\spring-boot-quiz-app-main')
    
    print("Creating comprehensive Quiz Platform presentation...")
    prs = create_presentation()
    
    output_file = "QUIZ_PLATFORM_PROFESSIONAL.pptx"
    prs.save(output_file)
    
    print(f"✅ Professional presentation created successfully!")
    print(f"📁 File: {output_file}")
    print(f"📍 Location: {os.path.abspath(output_file)}")
    print(f"📊 Total Slides: {len(prs.slides)}")
